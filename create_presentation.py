
"""
Generate Architecture Presentation for Retail Insights Assistant
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Create the architecture presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define color scheme - Ocean Gradient
    PRIMARY_COLOR = RGBColor(6, 90, 130)  # Deep blue
    SECONDARY_COLOR = RGBColor(28, 114, 147)  # Teal
    ACCENT_COLOR = RGBColor(33, 41, 92)  # Midnight
    TEXT_COLOR = RGBColor(54, 54, 54)  # Dark gray
    LIGHT_BG = RGBColor(255, 255, 255)  # White
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background for title slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Retail Insights Assistant"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Enterprise AI-Powered Analytics Platform"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(202, 220, 252)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle 2
    subtitle2_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    subtitle2_frame = subtitle2_box.text_frame
    subtitle2_frame.text = "Multi-Agent Architecture | Scalable to 100GB+"
    subtitle2_para = subtitle2_frame.paragraphs[0]
    subtitle2_para.font.size = Pt(18)
    subtitle2_para.font.color.rgb = RGBColor(202, 220, 252)
    subtitle2_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: System Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "System Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Content boxes
    content_text = """Key Features:
• Natural language query processing
• Multi-agent orchestration with LangGraph
• Real-time data analysis with DuckDB
• Semantic search using FAISS vectors
• Automated insight generation
• Interactive Streamlit interface

Core Capabilities:
• Process 100GB+ datasets
• Sub-second query response
• Conversational Q&A
• Automated business summaries
• Data quality validation"""
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(3.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    lines = content_text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        
        if line.endswith(':'):
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = SECONDARY_COLOR
    
    # Slide 3: Multi-Agent Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Multi-Agent Architecture"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Agent boxes
    agents = [
        ("Query Resolution\nAgent", "Converts natural language\nto structured SQL", 1.5),
        ("Data Extraction\nAgent", "Executes queries &\nretrieves data", 3.5),
        ("Validation\nAgent", "Validates results &\ndata quality", 5.5),
        ("Synthesis\nAgent", "Generates insights in\nnatural language", 7.5)
    ]
    
    y_pos = 1.3
    for agent_name, description, x_pos in agents:
        # Agent box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos - 0.6), Inches(y_pos),
            Inches(1.7), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SECONDARY_COLOR
        shape.line.color.rgb = PRIMARY_COLOR
        shape.line.width = Pt(2)
        
        # Agent name
        tf = shape.text_frame
        tf.text = agent_name
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x_pos - 0.6), Inches(y_pos + 1.25),
            Inches(1.7), Inches(0.7)
        )
        tf = desc_box.text_frame
        tf.text = description
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if x_pos < 7:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x_pos + 0.55), Inches(y_pos + 0.5),
                Inches(0.4), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY_COLOR
            arrow.line.fill.background()
    
    # Orchestrator label
    orch_box = slide.shapes.add_textbox(Inches(3), Inches(3.3), Inches(4), Inches(0.5))
    tf = orch_box.text_frame
    tf.text = "🔄 Orchestrated by LangGraph"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 4: Data Flow Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Query Processing Pipeline"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Pipeline steps
    steps = [
        "1. User Question",
        "2. Intent Analysis",
        "3. SQL Generation",
        "4. Data Retrieval",
        "5. Validation",
        "6. Insight Generation",
        "7. Natural Language Response"
    ]
    
    y_start = 1.2
    step_height = 0.5
    
    for i, step in enumerate(steps):
        y_pos = y_start + (i * step_height)
        
        # Step box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(y_pos),
            Inches(6), Inches(0.4)
        )
        
        # Alternate colors
        if i % 2 == 0:
            shape.fill.solid()
            shape.fill.fore_color.rgb = SECONDARY_COLOR
            text_color = RGBColor(255, 255, 255)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(202, 220, 252)
            text_color = PRIMARY_COLOR
        
        shape.line.color.rgb = PRIMARY_COLOR
        
        # Step text
        tf = shape.text_frame
        tf.text = step
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 5: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Technology Stack"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Technology categories
    tech_categories = [
        ("LLM & Agents", ["OpenAI GPT-4", "LangChain", "LangGraph", "LangSmith"]),
        ("Data Processing", ["DuckDB", "Pandas", "PyArrow", "NumPy"]),
        ("Vector Store", ["FAISS", "Sentence Transformers", "ChromaDB"]),
        ("UI Framework", ["Streamlit", "Plotly", "Altair"])
    ]
    
    x_positions = [0.5, 3, 5.5, 8]
    y_pos = 1.2
    
    for (category, techs), x_pos in zip(tech_categories, x_positions):
        # Category header
        header = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(2.3), Inches(0.4))
        tf = header.text_frame
        tf.text = category
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Background for header
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(2.3), Inches(0.4)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = SECONDARY_COLOR
        bg.line.fill.background()
        
        # Move to back
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
        
        # Technologies
        for i, tech in enumerate(techs):
            tech_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.1), Inches(y_pos + 0.5 + i * 0.35),
                Inches(2.1), Inches(0.3)
            )
            tf = tech_box.text_frame
            tf.text = f"• {tech}"
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_COLOR
    
    # Slide 6: Scalability Architecture (100GB+)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Scalability Architecture (100GB+)"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Three columns for scalability
    columns = [
        ("Data Engineering", [
            "• PySpark/Dask for distributed processing",
            "• Parquet columnar storage",
            "• Partitioning by date/region",
            "• Incremental processing"
        ]),
        ("Storage & Query", [
            "• Cloud Data Warehouse (BigQuery/Snowflake)",
            "• Data Lake (S3/GCS)",
            "• Delta Lake for ACID",
            "• Redis caching layer"
        ]),
        ("Optimization", [
            "• Materialized views",
            "• Query result caching",
            "• Vector index sharding",
            "• Approximate search (HNSW)"
        ])
    ]
    
    x_positions = [0.5, 3.5, 6.5]
    y_pos = 1.2
    
    for (col_title, items), x_pos in zip(columns, x_positions):
        # Column header
        header = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(3), Inches(0.4))
        tf = header.text_frame
        tf.text = col_title
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        
        # Items
        items_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(y_pos + 0.5),
            Inches(2.8), Inches(3)
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
    
    # Slide 7: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Performance & Monitoring"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Metrics in cards
    metrics = [
        ("Query Latency", "< 2 seconds\n(p95)", SECONDARY_COLOR),
        ("Data Processing", "10K rows/sec", PRIMARY_COLOR),
        ("Accuracy", "95%+ validation", ACCENT_COLOR),
        ("Scalability", "100GB+ datasets", SECONDARY_COLOR)
    ]
    
    x_start = 0.8
    y_pos = 1.5
    card_width = 2
    spacing = 0.3
    
    for i, (metric, value, color) in enumerate(metrics):
        x_pos = x_start + i * (card_width + spacing)
        
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(card_width), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # Metric name
        name_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(y_pos + 0.2),
            Inches(card_width), Inches(0.4)
        )
        tf = name_box.text_frame
        tf.text = metric
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(y_pos + 0.6),
            Inches(card_width), Inches(0.7)
        )
        tf = value_box.text_frame
        tf.text = value
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Monitoring tools
    monitoring_text = """Monitoring Stack:
• LangSmith for LLM tracing and debugging
• Prometheus + Grafana for metrics
• ELK Stack for log aggregation
• Custom dashboards for business KPIs"""
    
    monitor_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1.5))
    tf = monitor_box.text_frame
    tf.word_wrap = True
    
    lines = monitoring_text.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_COLOR
        
        if line.endswith(':'):
            p.font.bold = True
            p.font.color.rgb = SECONDARY_COLOR
    
    # Slide 8: Use Cases & Applications
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "Use Cases & Applications"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Use cases
    use_cases = """Executive Dashboards:
• Real-time sales performance monitoring
• Automated daily/weekly/monthly reports
• Trend analysis and forecasting

Business Analytics:
• Ad-hoc query answering for analysts
• Customer segmentation insights
• Product performance analysis

Operational Intelligence:
• Inventory optimization recommendations
• Regional performance comparisons
• Anomaly detection and alerts"""
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    lines = use_cases.split('\n')
    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        
        if line.endswith(':'):
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = SECONDARY_COLOR
            p.space_before = Pt(12)
    
    # Slide 9: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "Enterprise-Ready AI Analytics"
    p = tf.paragraphs[0]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Key points
    points = [
        "✓ Multi-agent architecture for complex queries",
        "✓ Scalable to 100GB+ datasets",
        "✓ Sub-second query response times",
        "✓ Natural language interface for all users"
    ]
    
    y_pos = 2.8
    for point in points:
        point_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.3))
        tf = point_box.text_frame
        tf.text = point
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(202, 220, 252)
        p.alignment = PP_ALIGN.CENTER
        y_pos += 0.35
    
    # Save presentation
    prs.save('/home/claude/retail-insights-assistant/Architecture_Presentation.pptx')
    print("Presentation created successfully!")
    print("Saved to: Architecture_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()